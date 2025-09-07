from io import BytesIO
from pptx import Presentation
from pptx.util import Pt
from pptx.dml.color import RGBColor

def hex_to_rgb(hex_color: str):
    h = hex_color.strip().lstrip("#")
    return tuple(int(h[i:i+2], 16) for i in (0, 2, 4))

def apply_brand_to_title(shape, brand_rgb):
    if shape.has_text_frame and shape.text_frame.paragraphs and shape.text_frame.paragraphs[0].runs:
        run = shape.text_frame.paragraphs[0].runs[0]
        run.font.color.rgb = RGBColor(*brand_rgb)

def add_title_slide(prs, title_text, brand_rgb):
    slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    if len(slide.placeholders) > 1:
        subtitle = slide.placeholders[1]
        subtitle.text = "Auto-generated deck"
    apply_brand_to_title(slide.shapes.title, brand_rgb)
    return slide

def add_content_slide(prs, title_text, bullets, notes, brand_rgb):
    slide_layout = prs.slide_layouts[1]
    slide = prs.slides.add_slide(slide_layout)
    slide.shapes.title.text = title_text
    apply_brand_to_title(slide.shapes.title, brand_rgb)

    body = slide.shapes.placeholders[1].text_frame
    body.clear()
    for i, b in enumerate(bullets):
        if i == 0 and body.paragraphs:
            p = body.paragraphs[0]
        else:
            p = body.add_paragraph()
        p.text = b
        p.level = 0
        for run in p.runs:
            run.font.size = Pt(20)

    if notes:
        slide.notes_slide.notes_text_frame.text = notes
    return slide

def deck_to_pptx_bytes(deck):
    prs = Presentation()
    brand_rgb = hex_to_rgb(deck.brand_color)

    add_title_slide(prs, deck.topic, brand_rgb)

    content_slides = deck.slides[1:] if deck.slides and deck.slides[0].bullets == [] else deck.slides
    for s in content_slides:
        title = s.title or "Slide"
        bullets = s.bullets or []
        notes = s.notes or ""
        add_content_slide(prs, title, bullets, notes, brand_rgb)

    bio = BytesIO()
    prs.save(bio)
    bio.seek(0)
    return bio.getvalue()
